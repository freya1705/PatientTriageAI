"""
Publication-Grade PowerPoint Generator for PatientTriage.ai Business Proposal
Generates PatientTriage_AI_Business_Proposal.pptx for Accenture Innovation Challenge 2026 Round 2
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    C_NAVY_DARK = RGBColor(15, 23, 42)      # #0f172a
    C_SLATE_CARD = RGBColor(30, 41, 59)     # #1e293b
    C_CYAN_ACCENT = RGBColor(2, 132, 199)   # #0284c7
    C_EMERALD = RGBColor(16, 185, 129)      # #10b981
    C_ROSE = RGBColor(225, 29, 72)          # #e11d48
    C_AMBER = RGBColor(245, 158, 11)        # #f59e0b
    C_WHITE = RGBColor(255, 255, 255)
    C_MUTED = RGBColor(148, 163, 184)       # #94a3b8
    C_LIGHT_BG = RGBColor(248, 250, 252)    # #f8fafc

    def add_bg(slide, dark=True):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_NAVY_DARK if dark else C_LIGHT_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="ACCENTURE INNOVATION CHALLENGE 2026", dark=True):
        # Category pill/subtitle
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(10)
        p_c.font.bold = True
        p_c.font.color.rgb = C_CYAN_ACCENT

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.6))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = C_WHITE if dark else C_NAVY_DARK

    def add_card(slide, left, top, width, height, bg_color=C_SLATE_CARD, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # ==========================================================
    # SLIDE 1: Title & Vision
    # ==========================================================
    s1 = prs.slides.add_slide(blank_layout)
    add_bg(s1, dark=True)

    # Main Hero Box
    hb = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.0))
    tf1 = hb.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "PATIENTTRIAGE.AI"
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = C_WHITE

    p1 = tf1.add_paragraph()
    p1.text = "“Triage is a snapshot. Risk isn't.”"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = C_CYAN_ACCENT
    p1.space_before = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "The Active Autonomous Emergency Department Safety Control Tower"
    p2.font.size = Pt(16)
    p2.font.color.rgb = C_MUTED
    p2.space_before = Pt(6)

    # 4 Feature Metric Cards at Bottom
    metrics = [
        ("CORE INNOVATION", "Dynamic Risk Velocity", "Vital shift tracking vs static snapshot", C_CYAN_ACCENT),
        ("PHILOSOPHY", "Unknown ≠ Safe", "Missing data penalizes confidence", C_ROSE),
        ("ANNUAL NET VALUE", "$3.82M / Hospital", "64% Avoided ICU Transfers + LWBS", C_EMERALD),
        ("AUTHOR & TRACK", "Freya Jadhav (freya1705)", "Round 2 Technical Prototype & Proposal", C_WHITE)
    ]
    card_w = Inches(2.7)
    card_h = Inches(1.8)
    for i, (kpi, val, desc, col) in enumerate(metrics):
        left = Inches(0.8 + i * 2.95)
        add_card(s1, left, Inches(4.8), card_w, card_h)
        tb = s1.shapes.add_textbox(left + Inches(0.15), Inches(4.9), card_w - Inches(0.3), card_h - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = kpi
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = C_MUTED
        
        pv = tf.add_paragraph()
        pv.text = val
        pv.font.size = Pt(14)
        pv.font.bold = True
        pv.font.color.rgb = col
        pv.space_before = Pt(4)
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9)
        pd.font.color.rgb = C_MUTED
        pd.space_before = Pt(3)

    # ==========================================================
    # SLIDE 2: The Problem & Clinical Crisis
    # ==========================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_bg(s2, dark=True)
    add_header(s2, "The Crisis: The Blind Spot in Emergency Waiting Rooms")

    p_cards = [
        ("1. Silent Decompensation", "140M annual visits face 3–6 hr wait times. ESI 3/4 patients worsen unmonitored; physiological decline may remain undetected until sudden collapse.", C_ROSE),
        ("2. The 'Missing Data Is Safe' Myth", "Traditional triage algorithms treat missing vitals as benign. Under 'Unknown is NOT Safe', incomplete data heightens vigilance instead of offering false reassurance.", C_AMBER),
        ("3. The Attention Bottleneck", "Attended critical patients anchor doctors and block queues, while deteriorating unmonitored patients stay hidden at the bottom of the list.", C_CYAN_ACCENT),
        ("4. Legacy EHR Blind Spot", "Epic EDI & Cerner MEWS only score admitted inpatient beds. PatientTriage.ai is purpose-built for waiting rooms and clinician coverage gaps.", C_EMERALD)
    ]
    for i, (title, desc, col) in enumerate(p_cards):
        row = i // 2
        col_idx = i % 2
        l = Inches(0.8 + col_idx * 5.95)
        t = Inches(1.7 + row * 2.6)
        add_card(s2, l, t, Inches(5.75), Inches(2.35), border_color=col)
        
        tb = s2.shapes.add_textbox(l + Inches(0.25), t + Inches(0.2), Inches(5.25), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = col
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = C_WHITE
        pd.space_before = Pt(8)

    # ==========================================================
    # SLIDE 3: System Architecture & The Attention Gap
    # ==========================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_bg(s3, dark=True)
    add_header(s3, "System Architecture: 3-Tier Layered Safety & The Attention Gap")

    tiers = [
        ("TIER 1: DETERMINISTIC GUARDRAILS", "• Hard clinical red-flags (SpO₂ < 85%, SBP < 75 mmHg, Stroke FAST, stridor)\n• 100% Downgrade Guardrails (prevent unsafe score lowering without evidence)", C_ROSE),
        ("TIER 2: AI TRAJECTORY & ATTENTION GAP", "• Priority = Base Risk + Vital Velocity (ΔV) + Evidence Staleness (τ) + Uncertainty (Unknown≠Safe) - Physician Coverage (wc)\n• 108 EMS Pre-Arrival Ingestion via HL7 FHIR LOINC\n• Referral Eligibility Scoring (RES 0–100%)", C_CYAN_ACCENT),
        ("TIER 3: MULTI-SURFACE DISPATCH & GOVERNANCE", "• 4 Clinical Workspaces: Overview, Nurse Tasks, Floor Map, Lab Orders\n• Unified slide-over Patient Drawer with human-readable 'Why?'\n• 100% Clinician override authority + Immutable append-only audit trail", C_EMERALD)
    ]
    for i, (t_title, t_body, col) in enumerate(tiers):
        l = Inches(0.8 + i * 3.95)
        add_card(s3, l, Inches(1.7), Inches(3.8), Inches(5.2), border_color=col)
        
        tb = s3.shapes.add_textbox(l + Inches(0.2), Inches(1.9), Inches(3.4), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = t_title
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = col
        
        pd = tf.add_paragraph()
        pd.text = t_body
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = C_WHITE
        pd.space_before = Pt(10)

    # ==========================================================
    # SLIDE 4: Clinical Command Cockpit (Live Demo Features)
    # ==========================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_bg(s4, dark=True)
    add_header(s4, "The Clinical Cockpit: 4 Switchable Workspaces")

    workspaces = [
        ("🧭 1. Overview (Control Tower)", "Dynamic Priority Stream continuously re-ranking patients by unmet clinical need + Searchable census + Live audit feed.", C_CYAN_ACCENT),
        ("🩺 2. Nurse Tasks ('Next 5 Mins')", "Time-budgeted micro-tasks (90s, 60s, 45s, 30s) converting abstract scores into direct bedside clinical actions + 'Attendant Away' spot-checks.", C_ROSE),
        ("🗺️ 3. ED Floor Pressure Map", "Real-time spatial visualization of Waiting Lounge Chairs 1–20 (pulsing halos 🟢/🟡/🟠/🔴) vs Treatment Bays + Unaccompanied patient alerts.", C_AMBER),
        ("🧪 4. Standing Lab Pre-Orders Hub", "Auto-drafts Troponin+ECG and Lactate orders before MD assignment with 1-click approval, cutting diagnostic turnaround time by 18 minutes.", C_EMERALD)
    ]
    for i, (w_title, w_desc, col) in enumerate(workspaces):
        row = i // 2
        col_idx = i % 2
        l = Inches(0.8 + col_idx * 5.95)
        t = Inches(1.7 + row * 2.6)
        add_card(s4, l, t, Inches(5.75), Inches(2.35), border_color=col)
        
        tb = s4.shapes.add_textbox(l + Inches(0.25), t + Inches(0.2), Inches(5.25), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = w_title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = col
        
        pd = tf.add_paragraph()
        pd.text = w_desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = C_WHITE
        pd.space_before = Pt(8)

    # ==========================================================
    # SLIDE 5: Patient Transparency Companion & De-Escalation
    # ==========================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_bg(s5, dark=True)
    add_header(s5, "Patient Transparency Companion: Reducing Anxiety & LWBS")

    s5_boxes = [
        ("📱 Zero-App Install (QR Code / SMS)", "Waiting patients scan a QR code printed on their triage wristband to access a live, mobile-friendly care journey tracker with zero app downloads.", C_CYAN_ACCENT),
        ("🤝 Behavioral De-Escalation ('Why Queue Moves')", "Directly answers: 'Why did someone who arrived after me get called in first?' — explaining clinical triage acuity without exposing private PHI.", C_EMERALD),
        ("⏱️ Next Nurse Re-Check Countdown", "Displays a reassuring countdown timer for the next scheduled vital check, eliminating waiting room uncertainty and panic.", C_AMBER),
        ("💰 $1.12M LWBS Revenue Recovery", "By keeping waiting patients informed and engaged, reduces Left-Without-Being-Seen (LWBS) walkout rates from 4.8% to <2.9%.", C_ROSE)
    ]
    for i, (title, desc, col) in enumerate(s5_boxes):
        row = i // 2
        col_idx = i % 2
        l = Inches(0.8 + col_idx * 5.95)
        t = Inches(1.7 + row * 2.6)
        add_card(s5, l, t, Inches(5.75), Inches(2.35), border_color=col)
        
        tb = s5.shapes.add_textbox(l + Inches(0.25), t + Inches(0.2), Inches(5.25), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = col
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = C_WHITE
        pd.space_before = Pt(8)

    # ==========================================================
    # SLIDE 6: Financial ROI Waterfall ($3.82M Net Value)
    # ==========================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_bg(s6, dark=True)
    add_header(s6, "Quantified Hospital Business Case & ROI ($3.82M Net Annual Value)")

    # 5 ROI Breakdown Columns
    roi_items = [
        ("1. AVOIDED ICU TRANSFERS", "+$1.39M", "64% reduction in waiting room crashes (93 avoided ICU stays @ $15k).", C_EMERALD),
        ("2. LWBS REVENUE RECOVERY", "+$1.12M", "30% reduction in walkouts via Patient Companion (936 patients retained).", C_CYAN_ACCENT),
        ("3. DIAGNOSTIC PRE-ORDERS", "+$830k", "18-min LOS reduction enables capacity for 660 additional admissions.", C_AMBER),
        ("4. MALPRACTICE MITIGATION", "+$480k", "40% liability reduction via tamper-evident audit ledger defense.", C_ROSE),
        ("TOTAL NET ANNUAL ROI", "+$3.58M", "Gross value $3.82M minus $240k enterprise license (14.9x Net ROI).", C_EMERALD)
    ]
    card_w6 = Inches(2.25)
    for i, (title, val, desc, col) in enumerate(roi_items):
        l = Inches(0.8 + i * 2.38)
        add_card(s6, l, Inches(1.7), card_w6, Inches(5.2), border_color=col)
        
        tb = s6.shapes.add_textbox(l + Inches(0.15), Inches(1.9), card_w6 - Inches(0.3), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(10)
        pt.font.bold = True
        pt.font.color.rgb = C_MUTED
        
        pv = tf.add_paragraph()
        pv.text = val
        pv.font.size = Pt(20)
        pv.font.bold = True
        pv.font.color.rgb = col
        pv.space_before = Pt(8)
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = C_WHITE
        pd.space_before = Pt(8)

    # ==========================================================
    # SLIDE 7: Regulatory, Roadmap & Conclusion
    # ==========================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_bg(s7, dark=True)
    add_header(s7, "Regulatory Compliance, Phased Roadmap & Outro")

    # Left: Regulatory & Privacy Box
    add_card(s7, Inches(0.8), Inches(1.7), Inches(5.75), Inches(5.2), border_color=C_CYAN_ACCENT)
    tb_l = s7.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.35), Inches(4.8))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "REGULATORY & ARCHITECTURE COMPLIANCE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_CYAN_ACCENT

    items_l = [
        ("FDA Non-Device CDS (21 U.S.C. § 360aaa-1)", "Transparent rationales & trajectories; clinicians maintain 100% final override authority."),
        ("Air-Gapped Edge Deployment", "Zero PHI leaves hospital network; sub-15ms inference latency on local server."),
        ("Immutable Audit Ledger", "Cryptographically verifiable tamper-evident SQLite WAL logging for all clinical actions."),
        ("Test Verification", "51/51 automated pytest tests passed (100% pass rate).")
    ]
    for k, v in items_l:
        p_k = tf_l.add_paragraph()
        p_k.text = f"• {k}: "
        p_k.font.size = Pt(10.5)
        p_k.font.bold = True
        p_k.font.color.rgb = C_WHITE
        p_k.space_before = Pt(8)
        
        p_v = tf_l.add_paragraph()
        p_v.text = f"  {v}"
        p_v.font.size = Pt(9.5)
        p_v.font.color.rgb = C_MUTED

    # Right: Phased Roadmap Box
    add_card(s7, Inches(6.78), Inches(1.7), Inches(5.75), Inches(5.2), border_color=C_EMERALD)
    tb_r = s7.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.35), Inches(4.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p_r = tf_r.paragraphs[0]
    p_r.text = "PHASED 12-MONTH EXECUTION ROADMAP"
    p_r.font.size = Pt(14)
    p_r.font.bold = True
    p_r.font.color.rgb = C_EMERALD

    phases = [
        ("Phase 1 (Months 1–3) — Edge Pilot", "Deploy contactless camera testbed in 2 triage bays; benchmark sensitivity across 51 test suites."),
        ("Phase 2 (Months 4–6) — Shadow Surveillance", "Integrate HL7 FHIR stream & 108 EMS pre-arrival telemetry; evaluate concordance alongside Epic/Cerner."),
        ("Phase 3 (Months 7–9) — Active Cockpit", "Launch Nurse Tasks ('Next 5 Mins') and Standing Pre-Orders with charge nurse feedback."),
        ("Phase 4 (Months 10–12) — Enterprise Scaling", "Roll out Patient Companion QR portal & Referral Diversion across 3 regional hospital networks.")
    ]
    for k, v in phases:
        p_pk = tf_r.add_paragraph()
        p_pk.text = f"• {k}"
        p_pk.font.size = Pt(10.5)
        p_pk.font.bold = True
        p_pk.font.color.rgb = C_WHITE
        p_pk.space_before = Pt(8)
        
        p_pv = tf_r.add_paragraph()
        p_pv.text = f"  {v}"
        p_pv.font.size = Pt(9.5)
        p_pv.font.color.rgb = C_MUTED

    # Save
    pptx_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "PatientTriage_AI_Business_Proposal.pptx")
    prs.save(pptx_path)
    print(f"[OK] Generated business proposal PPTX: {pptx_path} (Size: {os.path.getsize(pptx_path)} bytes)")

if __name__ == "__main__":
    create_deck()
